"""
PPTX 构建器 — 将解析后的 HTML 幻灯片数据重建为可编辑的 PPTX 文件。
借鉴 slide-gen 的 OOXML 重建思路，将 DOM 元素映射为原生 PowerPoint 对象。
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

from .html_parser import SlideData, SlideElement

# 16:9 幻灯片尺寸 (英寸)
SLIDE_WIDTH = 13.333
SLIDE_HEIGHT = 7.5

# 默认样式
DEFAULT_FONT = 'Microsoft YaHei'
FALLBACK_FONTS = ['Arial', 'Calibri', 'PingFang SC']
TITLE_SIZE = Pt(36)
HEADING_SIZE = Pt(24)
BODY_SIZE = Pt(16)
SMALL_SIZE = Pt(12)


def _hex_to_rgb(hex_color: str) -> RGBColor:
    """将 #RRGGBB 或 RRGGBB 转换为 RGBColor"""
    h = hex_color.lstrip('#')
    try:
        return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))
    except (ValueError, IndexError):
        return RGBColor(0x2D, 0x2D, 0x2D)


def _resolve_font(styles: dict) -> str:
    """从样式字典中解析 font-family"""
    font = styles.get('font-family', '')
    for f in FALLBACK_FONTS:
        if f.lower() in font.lower():
            return f
    return FALLBACK_FONTS[0]


def _resolve_font_size(styles: dict, default=Pt(16)) -> Pt:
    """从样式字典中解析 font-size"""
    fs = styles.get('font-size', '')
    if 'px' in fs:
        try:
            val = float(fs.replace('px', ''))
            return Pt(val * 0.75)  # px → pt
        except ValueError:
            pass
    elif 'pt' in fs:
        try:
            return Pt(float(fs.replace('pt', '')))
        except ValueError:
            pass
    return default


def _resolve_color(styles: dict, prop='color', default='#2d2d2d') -> RGBColor:
    """从样式字典中解析颜色"""
    c = styles.get(prop, default)
    if c.startswith('#'):
        return _hex_to_rgb(c)
    return _hex_to_rgb(default)


def _resolve_alignment(styles: dict) -> PP_ALIGN:
    """从样式字典中解析 text-align"""
    ta = styles.get('text-align', 'left')
    mapping = {
        'center': PP_ALIGN.CENTER,
        'right': PP_ALIGN.RIGHT,
        'left': PP_ALIGN.LEFT,
        'justify': PP_ALIGN.JUSTIFY,
    }
    return mapping.get(ta, PP_ALIGN.LEFT)


class PptxBuilder:
    """将 SlideData 列表构建为可编辑的 PPTX 文件"""

    def __init__(self, slide_width_inches=SLIDE_WIDTH,
                 slide_height_inches=SLIDE_HEIGHT):
        self.prs = Presentation()
        self.prs.slide_width = Inches(slide_width_inches)
        self.prs.slide_height = Inches(slide_height_inches)
        # 使用空白布局
        self.blank_layout = self.prs.slide_layouts[6]  # Blank

    def add_textbox(self, slide, left, top, width, height, text,
                    font_size=BODY_SIZE, color=None, bold=False,
                    alignment=PP_ALIGN.LEFT, font_name=None):
        """添加可编辑文本框"""
        txBox = slide.shapes.add_textbox(
            Inches(left), Inches(top), Inches(width), Inches(height)
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = font_size
        p.font.color.rgb = color or _hex_to_rgb('#2d2d2d')
        p.font.bold = bold
        p.font.name = font_name or DEFAULT_FONT
        p.alignment = alignment
        return txBox

    def add_bullet_list(self, slide, left, top, width, height,
                        items: list[str], font_size=BODY_SIZE,
                        color=None):
        """添加可编辑的要点列表"""
        txBox = slide.shapes.add_textbox(
            Inches(left), Inches(top), Inches(width), Inches(height)
        )
        tf = txBox.text_frame
        tf.word_wrap = True

        for i, item in enumerate(items):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = item
            p.font.size = font_size
            p.font.color.rgb = color or _hex_to_rgb('#555555')
            p.font.name = DEFAULT_FONT
            p.space_after = Pt(6)
            # 要点符号
            p.level = 0
            pPr = p._pPr
            if pPr is None:
                from pptx.oxml.ns import qn
                pPr = p._p.get_or_add_pPr()
            from pptx.oxml.ns import qn
            from lxml import etree
            buChar = etree.SubElement(pPr, qn('a:buChar'))
            buChar.set('char', '•')
        return txBox

    def add_image_shape(self, slide, left, top, width, height,
                        image_path: str):
        """添加图片"""
        if Path(image_path).exists():
            slide.shapes.add_picture(
                image_path,
                Inches(left), Inches(top),
                Inches(width), Inches(height)
            )

    def add_rectangle(self, slide, left, top, width, height,
                      fill_color=None, border_color=None):
        """添加矩形形状（作为装饰/背景）"""
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(left), Inches(top),
            Inches(width), Inches(height)
        )
        if fill_color:
            shape.fill.solid()
            shape.fill.fore_color.rgb = _hex_to_rgb(fill_color)
        else:
            shape.fill.background()

        if border_color:
            shape.line.color.rgb = _hex_to_rgb(border_color)
            shape.line.width = Pt(1)
        else:
            shape.line.fill.background()
        return shape

    def build_slide_from_data(self, slide_data: SlideData,
                              screenshot_path: str = None):
        """根据 SlideData 构建一页幻灯片"""
        slide = self.prs.slides.add_slide(self.blank_layout)

        # 提取标题
        title = slide_data.title or f'第 {slide_data.index + 1} 页'

        # 如果有截图，优先使用截图作为背景（保真度最高）
        if screenshot_path and Path(screenshot_path).exists():
            slide.shapes.add_picture(
                screenshot_path,
                Inches(0), Inches(0),
                Inches(SLIDE_WIDTH), Inches(SLIDE_HEIGHT)
            )
            return slide

        # 否则: 构建可编辑元素
        # 添加标题
        title_color = _hex_to_rgb('#E8695A')
        self.add_textbox(
            slide, 0.8, 0.4, 11.7, 0.9,
            title, font_size=TITLE_SIZE, color=title_color,
            bold=True
        )

        # 添加分隔线
        self.add_rectangle(slide, 0.8, 1.25, 2.0, 0.04,
                           fill_color='#E8695A')

        # 添加元素（跳过已渲染的标题，避免重复）
        rendered_titles = {title.strip()}
        y_offset = 1.6
        for el in slide_data.elements:
            if y_offset > 6.5:
                break
            if el.tag in ['h1', 'h2', 'h3', 'h4']:
                # 跳过已在顶部渲染的标题
                if el.text.strip() in rendered_titles:
                    continue
                rendered_titles.add(el.text.strip())
                fs = _resolve_font_size(el.styles, HEADING_SIZE)
                color = _resolve_color(el.styles)
                self.add_textbox(
                    slide, 1.2, y_offset, 10.5, 0.6,
                    el.text, font_size=fs, color=color, bold=True
                )
                y_offset += 0.55
            elif el.tag in ['p', 'span', 'div'] and el.text:
                fs = _resolve_font_size(el.styles, BODY_SIZE)
                color = _resolve_color(el.styles)
                align = _resolve_alignment(el.styles)
                lines = el.text.split('\n')
                height = max(0.35, len(lines) * 0.3)
                self.add_textbox(
                    slide, 1.2, y_offset, 10.5, height,
                    el.text, font_size=fs, color=color,
                    alignment=align
                )
                y_offset += height + 0.1
            elif el.tag in ['ul', 'ol']:
                # 列表项处理
                items = el.text.split('\n')
                items = [i.strip() for i in items if i.strip()]
                if items:
                    height = max(0.5, len(items) * 0.3)
                    self.add_bullet_list(
                        slide, 1.2, y_offset, 10.5, height,
                        items, font_size=BODY_SIZE
                    )
                    y_offset += height + 0.1

        return slide

    def save(self, output_path: str):
        """保存 PPTX 文件"""
        path = Path(output_path)
        path.parent.mkdir(parents=True, exist_ok=True)
        self.prs.save(str(path))
        return str(path)
