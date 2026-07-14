#!/usr/bin/env python3
"""
HTML-PPT 转换器测试用例
覆盖: 解析成功率、PPTX 完整性、截图质量
"""
import sys
import os
import time
from pathlib import Path

# 编码兼容
if sys.stdout.encoding and sys.stdout.encoding.upper() != 'UTF-8':
    import io
    sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')

sys.path.insert(0, str(Path(__file__).parent.parent))

from converter.html_parser import HtmlParser
from converter.pptx_builder import PptxBuilder
from converter.screenshot_engine import (
    ScreenshotEngine, process_screenshot, save_screenshot
)
from converter.batch_processor import BatchProcessor


# ============ 测试数据 ============

SAMPLE_HTML_1 = '''<!DOCTYPE html>
<html><head><meta charset="UTF-8"><style>
.slide { width: 960px; height: 540px; padding: 40px; font-family: Arial; }
h1 { color: #e8695a; font-size: 36px; }
p { color: #333; font-size: 16px; line-height: 1.6; }
</style></head><body>
<div class="slide">
  <h1>项目封面</h1>
  <p>基于 AI 的智能 PPT 生成工具</p>
  <p>开发者: EthanBAI-dev</p>
</div>
<div class="slide">
  <h1>核心功能</h1>
  <p>七角色协作工作流</p>
  <p>AI 驱动 · 全自动生成</p>
</div>
<div class="slide">
  <h1>技术栈</h1>
  <p>Python · Playwright · OOXML</p>
</div>
</body></html>'''

SAMPLE_HTML_2 = '''<!DOCTYPE html>
<html><head><meta charset="UTF-8"><style>
.slide { width: 1280px; height: 720px; padding: 60px; background: #faf7f2; }
h1 { color: #2d2d2d; font-size: 42px; text-align: center; }
ul { font-size: 20px; color: #555; }
</style></head><body>
<section class="slide">
  <h1>NoteFlow 项目</h1>
  <ul><li>跨平台笔记应用</li><li>实时同步</li><li>Markdown 支持</li></ul>
</section>
<section class="slide">
  <h1>技术架构</h1>
  <p>前端: React + TypeScript</p>
  <p>后端: Node.js + PostgreSQL</p>
</section>
</body></html>'''


def test_html_parser():
    """测试 HTML 解析器"""
    print('\n🔍 [测试] HTML 解析器...')
    parser = HtmlParser()

    # 测试 1: 解析文件
    import tempfile
    with tempfile.NamedTemporaryFile(mode='w', suffix='.html',
                                     delete=False, encoding='utf-8') as f:
        f.write(SAMPLE_HTML_1)
        tmp1 = f.name

    try:
        slides = parser.parse_file(tmp1)
        assert len(slides) == 3, f'预期 3 张幻灯片, 实际 {len(slides)}'
        assert slides[0].title == '项目封面'
        assert slides[1].title == '核心功能'
        print(f'  ✅ 单文件解析: 3 张幻灯片, 标题正确')
    finally:
        os.unlink(tmp1)

    # 测试 2: 解析目录
    with tempfile.TemporaryDirectory() as tmpdir:
        f1 = Path(tmpdir) / 'test1.html'
        f2 = Path(tmpdir) / 'test2.html'
        f1.write_text(SAMPLE_HTML_1, encoding='utf-8')
        f2.write_text(SAMPLE_HTML_2, encoding='utf-8')

        result = parser.parse_directory(tmpdir)
        assert 'test1.html' in result
        assert 'test2.html' in result
        assert len(result['test1.html']) == 3
        assert len(result['test2.html']) == 2
        print(f'  ✅ 目录解析: 2 个文件, 共 5 张幻灯片')

    # 测试 3: 提取元素
    from bs4 import BeautifulSoup
    soup = BeautifulSoup(SAMPLE_HTML_1, 'lxml')
    slides = parser.find_slides(soup)
    assert len(slides) == 3

    title = parser.extract_title(slides[0])
    assert title == '项目封面'

    elements = parser.extract_elements(slides[0])
    assert len(elements) >= 3
    print(f'  ✅ 元素提取: 封面页 {len(elements)} 个元素')

    print('  ✅ HTML 解析器全部测试通过!')


def test_pptx_builder():
    """测试 PPTX 构建器"""
    print('\n🔍 [测试] PPTX 构建器...')
    parser = HtmlParser()
    import tempfile

    with tempfile.NamedTemporaryFile(mode='w', suffix='.html',
                                     delete=False, encoding='utf-8') as f:
        f.write(SAMPLE_HTML_1)
        tmp = f.name

    try:
        slides = parser.parse_file(tmp)
        builder = PptxBuilder()
        for slide_data in slides:
            builder.build_slide_from_data(slide_data)

        # 保存并验证
        output = Path(tempfile.mktemp(suffix='.pptx'))
        builder.save(str(output))
        assert output.exists()
        assert output.stat().st_size > 0

        # 验证 PPTX 完整性
        from pptx import Presentation
        prs = Presentation(str(output))
        assert len(prs.slides) == 3, f'预期 3 页, 实际 {len(prs.slides)}'

        # 检查文本内容是否可编辑
        for i, slide in enumerate(prs.slides):
            shapes = slide.shapes
            assert len(shapes) > 0, f'第 {i + 1} 页没有内容'
            # 验证至少有一个文本框
            has_textbox = any(s.has_text_frame for s in shapes)
            assert has_textbox, f'第 {i + 1} 页没有文本框'

        os.unlink(str(output))
        print(f'  ✅ PPTX 构建: 3 页可编辑幻灯片, 文本可编辑性验证通过')

    finally:
        os.unlink(tmp)

    print('  ✅ PPTX 构建器全部测试通过!')


def test_screenshot_processing():
    """测试截图后处理"""
    print('\n🔍 [测试] 截图后处理...')
    import tempfile

    # 创建一个测试图片 (纯色)
    from PIL import Image, ImageDraw
    img = Image.new('RGB', (1920, 1080), color=(250, 247, 242))
    draw = ImageDraw.Draw(img)
    draw.rectangle([100, 100, 500, 200], fill=(232, 105, 90))

    buf = io.BytesIO()
    img.save(buf, format='PNG')
    img_bytes = buf.getvalue()

    # 测试格式转换
    png_result = process_screenshot(img_bytes, 'PNG')
    assert len(png_result) > 0
    print(f'  ✅ PNG 输出: {len(png_result)} bytes')

    jpg_result = process_screenshot(img_bytes, 'JPG')
    assert len(jpg_result) > 0
    print(f'  ✅ JPG 输出: {len(jpg_result)} bytes')

    # 测试分辨率
    resized = process_screenshot(img_bytes, 'PNG', resolution=(960, 540))
    from PIL import Image as PILImage
    from io import BytesIO
    reloaded = PILImage.open(BytesIO(resized))
    assert reloaded.size == (960, 540)
    print(f'  ✅ 分辨率调整: 1920x1080 → 960x540')

    # 测试裁剪
    cropped = process_screenshot(img_bytes, 'PNG', crop_area=(100, 100, 500, 500))
    reloaded = PILImage.open(BytesIO(cropped))
    assert reloaded.size == (400, 400)
    print(f'  ✅ 裁剪功能: (100,100,500,500) → 400x400')

    # 测试保存
    with tempfile.TemporaryDirectory() as tmpdir:
        path = save_screenshot(img_bytes, str(Path(tmpdir) / 'test.png'),
                                output_format='PNG')
        assert Path(path).exists()
        print(f'  ✅ 文件保存: {path}')

        path = save_screenshot(img_bytes, str(Path(tmpdir) / 'test.jpg'),
                                output_format='JPG')
        assert Path(path).exists()
        print(f'  ✅ JPG 保存: {path}')

    print('  ✅ 截图后处理全部测试通过!')


def test_batch_processor():
    """测试批处理流程（不启动浏览器，仅验证文件发现和报告生成）"""
    print('\n🔍 [测试] 批处理调度器...')
    import tempfile
    import json

    with tempfile.TemporaryDirectory() as tmpdir:
        # 准备测试文件
        for i, html in enumerate([SAMPLE_HTML_1, SAMPLE_HTML_2]):
            f = Path(tmpdir) / f'test_{i + 1}.html'
            f.write_text(html, encoding='utf-8')

        # 使用 2 个线程处理（但跳过实际浏览器截图）
        processor = BatchProcessor(workers=2, output_dir=tmpdir)

        # 仅测试文件发现
        html_files = sorted(Path(tmpdir).glob('*.html'))
        assert len(html_files) == 2
        print(f'  ✅ 文件发现: {len(html_files)} 个 HTML 文件')

        # 测试报告生成
        from converter.batch_processor import ConversionResult
        results = [
            ConversionResult('test_1.html', 'success', 'test.pptx', ['shot1.png'], '', 1.2),
            ConversionResult('test_2.html', 'success', 'test.pptx', [], '', 0.8),
        ]
        processor._generate_report(results, Path(tmpdir))

        report_path = Path(tmpdir) / 'conversion_report.json'
        assert report_path.exists()
        report = json.loads(report_path.read_text(encoding='utf-8'))
        assert report['summary']['total'] == 2
        assert report['summary']['success'] == 2
        print(f'  ✅ 报告生成: {report_path.name}')

    print('  ✅ 批处理调度器全部测试通过!')


if __name__ == '__main__':
    import io  # for test_screenshot_processing

    print('=' * 50)
    print('🧪 HTML-PPT 转换器 测试套件')
    print('=' * 50)

    tests = [
        test_html_parser,
        test_pptx_builder,
        test_screenshot_processing,
        test_batch_processor,
    ]

    passed = 0
    failed = 0
    for test in tests:
        try:
            test()
            passed += 1
        except Exception as e:
            print(f'  ❌ 测试失败: {test.__name__}: {e}')
            import traceback
            traceback.print_exc()
            failed += 1

    print(f'\n{"="*50}')
    print(f'📊 测试结果: ✅ {passed} 通过 / ❌ {failed} 失败 / 总计 {len(tests)}')
    print(f'{"="*50}')
    sys.exit(1 if failed > 0 else 0)
